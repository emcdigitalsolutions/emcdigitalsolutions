"""
Genera proposta-commerciale.pptx — EMC Digital Solutions
9 slide con palette navy + blu + oro, coerente con la proposta HTML.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Palette ───
BG_PRIMARY = RGBColor(0x0C, 0x11, 0x20)
BG_SECONDARY = RGBColor(0x0A, 0x0E, 0x1A)
BG_CARD = RGBColor(0x17, 0x20, 0x35)
BG_TERTIARY = RGBColor(0x13, 0x1B, 0x2E)
ACCENT = RGBColor(0x6B, 0x9E, 0xF7)
ACCENT_LIGHT = RGBColor(0x93, 0xB8, 0xFC)
ACCENT_VIVID = RGBColor(0xF5, 0xC5, 0x42)
TEXT_PRIMARY = RGBColor(0xE8, 0xE4, 0xDC)
TEXT_SECONDARY = RGBColor(0x9C, 0xA3, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x0C, 0x11, 0x20)

FONT_HEADING = "Playfair Display"
FONT_BODY = "Inter"

# Slide dimensions: 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=18, color=TEXT_PRIMARY, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txbox


def add_multiline_textbox(slide, left, top, width, height, lines, font_name=FONT_BODY,
                          font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.LEFT,
                          line_spacing=1.4):
    """lines: list of (text, {overrides}) or just str"""
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, overrides = line, {}
        else:
            txt, overrides = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.name = overrides.get("font_name", font_name)
        p.font.size = Pt(overrides.get("font_size", font_size))
        p.font.color.rgb = overrides.get("color", color)
        p.font.bold = overrides.get("bold", False)
        p.font.italic = overrides.get("italic", False)
        p.alignment = overrides.get("alignment", alignment)
        p.space_after = Pt(overrides.get("space_after", 4))
        p.line_spacing = Pt(overrides.get("font_size", font_size) * line_spacing)
    return txbox


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_accent_line(slide, left, top, width):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return line


def add_stat_card(slide, left, top, number, label, width=2.8):
    add_rounded_rect(slide, left, top, width, 1.4, BG_CARD, ACCENT)
    add_textbox(slide, left, top + 0.15, width, 0.7, number,
                font_name=FONT_HEADING, font_size=32, color=ACCENT_LIGHT,
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.15, top + 0.75, width - 0.3, 0.6, label,
                font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


def add_pillar_card(slide, left, top, num, title, desc, width=3.4):
    add_rounded_rect(slide, left, top, width, 2.6, BG_CARD, ACCENT)
    add_textbox(slide, left, top + 0.2, width, 0.4, num,
                font_name=FONT_HEADING, font_size=12, color=ACCENT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + 0.6, width, 0.5, title,
                font_name=FONT_HEADING, font_size=18, color=ACCENT_LIGHT,
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.3, top + 1.2, width - 0.6, 1.2, desc,
                font_size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


def add_feature_card(slide, left, top, title, desc, width=5.2, height=1.2):
    add_rounded_rect(slide, left, top, width, height, BG_CARD, None)
    add_textbox(slide, left + 0.25, top + 0.1, width - 0.5, 0.4, title,
                font_size=14, color=TEXT_PRIMARY, bold=True)
    add_textbox(slide, left + 0.25, top + 0.5, width - 0.5, 0.7, desc,
                font_size=11, color=TEXT_SECONDARY)


def add_price_card(slide, left, top, plan_name, setup, monthly, features, featured=False):
    card_w = 3.6
    card_h = 5.0
    border = ACCENT_VIVID if featured else ACCENT
    add_rounded_rect(slide, left, top, card_w, card_h, BG_CARD, border)

    if featured:
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + card_w - 1.6), Inches(top + 0.15), Inches(1.4), Inches(0.32)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_VIVID
        badge.line.fill.background()
        badge.adjustments[0] = 0.3
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].text = "CONSIGLIATO"
        badge_tf.paragraphs[0].font.size = Pt(8)
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].font.color.rgb = DARK_TEXT
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_tf.paragraphs[0].font.name = FONT_BODY

    # Plan name
    add_textbox(slide, left + 0.3, top + 0.2, card_w - 0.6, 0.4, plan_name,
                font_name=FONT_HEADING, font_size=20, color=TEXT_PRIMARY, bold=True)
    # Setup
    add_textbox(slide, left + 0.3, top + 0.65, card_w - 0.6, 0.3, f"Setup: {setup}",
                font_size=12, color=TEXT_SECONDARY)
    # Monthly
    add_textbox(slide, left + 0.3, top + 1.0, card_w - 0.6, 0.5, monthly,
                font_name=FONT_HEADING, font_size=28, color=TEXT_PRIMARY, bold=True)

    # Separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + 0.3), Inches(top + 1.55), Inches(card_w - 0.6), Pt(1)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT if not featured else ACCENT_VIVID
    sep.line.fill.background()

    # Features list
    y = top + 1.75
    for feat in features:
        disabled = feat.startswith("~")
        text = feat.lstrip("~")
        bullet_color = TEXT_SECONDARY if disabled else (ACCENT_VIVID if featured else ACCENT)
        text_color = RGBColor(0x55, 0x5A, 0x65) if disabled else TEXT_SECONDARY

        # Bullet dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + 0.4), Inches(y + 0.08), Pt(6), Pt(6)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = bullet_color
        dot.line.fill.background()

        add_textbox(slide, left + 0.65, y - 0.05, card_w - 1.0, 0.35, text,
                    font_size=11, color=text_color)
        y += 0.38


# ═══════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank layout

# ════════════════════════════════════════════
# SLIDE 1: Copertina
# ════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, BG_SECONDARY)

# EMC badge
add_textbox(slide1, 4.5, 1.2, 4.3, 0.5, "EMC  DIGITAL SOLUTIONS",
            font_name=FONT_BODY, font_size=14, color=ACCENT_LIGHT,
            bold=True, alignment=PP_ALIGN.CENTER)

# Title
add_textbox(slide1, 1.5, 2.3, 10.3, 1.2, "Proposta di Visibilità Digitale",
            font_name=FONT_HEADING, font_size=44, color=TEXT_PRIMARY,
            bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide1, 3, 3.6, 7.3, 0.7, "Per la Tua Attività",
            font_name=FONT_HEADING, font_size=24, color=ACCENT,
            italic=True, alignment=PP_ALIGN.CENTER)

# Accent line
add_accent_line(slide1, 6.0, 4.5, 1.3)

# Date
add_textbox(slide1, 4.5, 4.9, 4.3, 0.5, "APRILE 2026",
            font_size=13, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER)

# Footer EMC
add_textbox(slide1, 3.5, 6.6, 6.3, 0.4,
            "Progettato e Sviluppato da EMC Digital Solutions",
            font_size=10, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 2: Il Problema
# ════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, BG_PRIMARY)

add_textbox(slide2, 0.8, 0.5, 5, 0.6, "Il Problema",
            font_name=FONT_HEADING, font_size=32, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide2, 0.8, 1.1, 1.0)

# Quote block
add_rounded_rect(slide2, 0.8, 1.5, 11.6, 1.3, BG_TERTIARY, ACCENT)
# Quote accent bar (left)
bar = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.5), Pt(5), Inches(1.3)
)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

add_textbox(slide2, 1.2, 1.6, 10.8, 0.7,
            '"Ho cercato un\'attività come la tua su Google... ma non ti ho trovato. Così ho scelto un altro."',
            font_name=FONT_HEADING, font_size=18, color=TEXT_PRIMARY, italic=True)
add_textbox(slide2, 1.2, 2.25, 10.8, 0.4,
            "— Ciò che i clienti pensano, ma non ti dicono",
            font_size=12, color=ACCENT)

# Insight box
add_rounded_rect(slide2, 0.8, 3.1, 11.6, 1.3, BG_TERTIARY, None)
add_textbox(slide2, 1.1, 3.2, 5, 0.4, 'Il concetto di "Top-of-Mind"',
            font_name=FONT_HEADING, font_size=16, color=ACCENT_LIGHT, bold=True)
add_textbox(slide2, 1.1, 3.6, 11.0, 0.7,
            "Quando una persona cerca un prodotto o un servizio, si rivolge al primo nome che trova. "
            "Non al migliore, non al più economico — al più visibile. Chi non è presente online, "
            "semplicemente non esiste.",
            font_size=13, color=TEXT_SECONDARY)

# Stats
add_stat_card(slide2, 1.2, 4.8, "87%", "degli italiani usa\ni social ogni giorno", 3.2)
add_stat_card(slide2, 5.0, 4.8, "72%", "cerca attività locali\nsu Google e Facebook", 3.2)
add_stat_card(slide2, 8.8, 4.8, "3 su 4", "scelgono l'attività che\ntrovano per prima online", 3.2)

# ════════════════════════════════════════════
# SLIDE 3: La Soluzione
# ════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, BG_PRIMARY)

add_textbox(slide3, 0.8, 0.5, 5, 0.6, "La Soluzione",
            font_name=FONT_HEADING, font_size=32, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide3, 0.8, 1.1, 1.0)

add_textbox(slide3, 0.8, 1.4, 9, 0.6,
            "Una strategia digitale completa costruita su tre pilastri, pensata per posizionare la tua attività come punto di riferimento nella zona.",
            font_size=15, color=TEXT_SECONDARY)

add_pillar_card(slide3, 0.8, 2.4, "01", "Presenza Online",
                "Sito web professionale che lavora per te 24 ore su 24, 7 giorni su 7.")
add_pillar_card(slide3, 4.7, 2.4, "02", "Social Media",
                "Pagine Facebook e Instagram curate con contenuti professionali e costanti.")
add_pillar_card(slide3, 8.6, 2.4, "03", "Visibilità Locale",
                "Google Business Profile e campagne sponsorizzate geolocalizzate nella tua area.")

# ════════════════════════════════════════════
# SLIDE 4: Sito Web (Servizio 01)
# ════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, BG_PRIMARY)

add_textbox(slide4, 0.8, 0.35, 1.2, 0.7, "01",
            font_name=FONT_HEADING, font_size=40, color=ACCENT, bold=True)
add_textbox(slide4, 2.0, 0.5, 6, 0.6, "Sito Web Professionale",
            font_name=FONT_HEADING, font_size=28, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide4, 2.0, 1.1, 1.0)

add_textbox(slide4, 0.8, 1.4, 10, 0.6,
            "La tua vetrina digitale aperta 24 ore su 24. Un punto di riferimento per chiunque cerchi i tuoi servizi online.",
            font_size=15, color=TEXT_SECONDARY)

add_feature_card(slide4, 0.8, 2.3, "Design su Misura",
                 "Ogni sito viene progettato su misura per la tua attività. Nessun template generico: solo soluzioni personalizzate.", 5.4)
add_feature_card(slide4, 6.6, 2.3, "Mobile-Friendly",
                 "Perfetto su smartphone e tablet. Il 70% delle visite arriva da dispositivi mobili.", 5.4)
add_feature_card(slide4, 0.8, 3.8, "SEO Ottimizzato",
                 "Posizionamento su Google per le ricerche più rilevanti nella tua zona e nel tuo settore.", 5.4)
add_feature_card(slide4, 6.6, 3.8, "Veloce e Sicuro",
                 "Tempi di caricamento rapidi, certificato SSL e hosting affidabile. Performance e sicurezza garantite.", 5.4)

# ════════════════════════════════════════════
# SLIDE 5: Social Media (Servizio 02)
# ════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, BG_PRIMARY)

add_textbox(slide5, 0.8, 0.35, 1.2, 0.7, "02",
            font_name=FONT_HEADING, font_size=40, color=ACCENT, bold=True)
add_textbox(slide5, 2.0, 0.5, 6, 0.6, "Gestione Social Media",
            font_name=FONT_HEADING, font_size=28, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide5, 2.0, 1.1, 1.0)

add_textbox(slide5, 0.8, 1.4, 10, 0.6,
            "Contenuti professionali e costanti che mantengono il tuo nome visibile nella mente delle persone, ogni giorno.",
            font_size=15, color=TEXT_SECONDARY)

# Facebook column
add_rounded_rect(slide5, 0.8, 2.2, 5.4, 3.2, BG_CARD, None)
add_textbox(slide5, 1.1, 2.3, 4, 0.4, "Facebook",
            font_name=FONT_HEADING, font_size=18, color=ACCENT_LIGHT, bold=True)
fb_items = [
    "Pagina Business ottimizzata",
    "8-12 post al mese",
    "Gestione commenti e messaggi",
    "Informazioni sempre aggiornate",
    "Orari, servizi e contatti visibili"
]
add_multiline_textbox(slide5, 1.1, 2.8, 4.8, 2.4,
                      [("›  " + item, {"font_size": 13}) for item in fb_items],
                      color=TEXT_SECONDARY, line_spacing=1.8)

# Instagram column
add_rounded_rect(slide5, 6.6, 2.2, 5.4, 3.2, BG_CARD, None)
add_textbox(slide5, 6.9, 2.3, 4, 0.4, "Instagram",
            font_name=FONT_HEADING, font_size=18, color=ACCENT_LIGHT, bold=True)
ig_items = [
    "Profilo curato e professionale",
    "Storie e contenuti visivi",
    "Immagini di qualità",
    "Hashtag strategici locali",
    "Engagement con la community"
]
add_multiline_textbox(slide5, 6.9, 2.8, 4.8, 2.4,
                      [("›  " + item, {"font_size": 13}) for item in ig_items],
                      color=TEXT_SECONDARY, line_spacing=1.8)

# Insight box
add_rounded_rect(slide5, 0.8, 5.6, 11.2, 1.1, BG_TERTIARY, None)
add_textbox(slide5, 1.1, 5.65, 5, 0.3, "Tipologie di contenuto",
            font_name=FONT_HEADING, font_size=14, color=ACCENT_LIGHT, bold=True)
add_textbox(slide5, 1.1, 5.95, 10.6, 0.7,
            "I tuoi prodotti e servizi, il dietro le quinte della tua attività, novità e promozioni, recensioni dei clienti, eventi e ricorrenze.",
            font_size=12, color=TEXT_SECONDARY)

# ════════════════════════════════════════════
# SLIDE 6: Sponsorizzate (Servizio 03)
# ════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6, BG_PRIMARY)

add_textbox(slide6, 0.8, 0.35, 1.2, 0.7, "03",
            font_name=FONT_HEADING, font_size=40, color=ACCENT, bold=True)
add_textbox(slide6, 2.0, 0.5, 6, 0.6, "Sponsorizzate Locali",
            font_name=FONT_HEADING, font_size=28, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide6, 2.0, 1.1, 1.0)

add_textbox(slide6, 0.8, 1.4, 10, 0.6,
            "Campagne mirate che portano il tuo nome davanti alle persone giuste, nel territorio giusto, al momento giusto.",
            font_size=15, color=TEXT_SECONDARY)

add_textbox(slide6, 0.8, 2.1, 5, 0.4, "Copertura Geografica",
            font_name=FONT_HEADING, font_size=16, color=ACCENT_LIGHT, bold=True)

# Geo tags
geo_labels = ["La tua città", "Comuni limitrofi", "La tua provincia", "Raggio personalizzabile"]
x_pos = 0.8
for label in geo_labels:
    tag_w = len(label) * 0.12 + 0.6
    tag = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_pos), Inches(2.55), Inches(tag_w), Inches(0.38)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = BG_TERTIARY
    tag.line.color.rgb = ACCENT
    tag.line.width = Pt(1)
    tag.adjustments[0] = 0.4
    tf = tag.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = ACCENT_LIGHT
    tf.paragraphs[0].font.name = FONT_BODY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_pos += tag_w + 0.2

# Checklist
check_items = [
    "Campagne geolocalizzate su Facebook e Instagram",
    "Budget contenuto con risultati misurabili",
    "Targeting per età, interessi e area geografica",
    "Report periodico con tutte le metriche"
]
add_multiline_textbox(slide6, 0.8, 3.2, 11, 1.8,
                      [("  " + item, {"font_size": 14}) for item in check_items],
                      color=TEXT_SECONDARY, line_spacing=1.8)

# Metric cards
add_stat_card(slide6, 1.2, 5.2, "5.000+", "Persone raggiunte\nal mese (stima)", 3.0)
add_stat_card(slide6, 5.0, 5.2, "da €0,02", "Costo medio\nper persona raggiunta", 3.0)
add_stat_card(slide6, 8.8, 5.2, "20 km", "Raggio di copertura\npersonalizzabile", 3.0)

# ════════════════════════════════════════════
# SLIDE 7: Investimento
# ════════════════════════════════════════════
slide7 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide7, BG_PRIMARY)

add_textbox(slide7, 0.8, 0.3, 5, 0.6, "Investimento",
            font_name=FONT_HEADING, font_size=32, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide7, 0.8, 0.88, 1.0)

add_textbox(slide7, 0.8, 1.1, 10, 0.5,
            "Pricing chiaro e trasparente, senza costi nascosti. Scegli il piano più adatto alla tua attività.",
            font_size=14, color=TEXT_SECONDARY)

# 3 Price Cards
add_price_card(slide7, 0.6, 1.7, "Base", "€300", "€150/mese", [
    "Sito web one-page",
    "1 piattaforma social",
    "4-6 post al mese",
    "~Sponsorizzate non incluse",
    "Report trimestrale",
    "~Google Business non incluso"
])

add_price_card(slide7, 4.7, 1.7, "Standard", "€500", "€200/mese", [
    "Sito web completo",
    "Facebook + Instagram",
    "8-12 post/mese + storie",
    "Gestione sponsorizzate inclusa",
    "Report mensile",
    "~Google Business non incluso"
], featured=True)

add_price_card(slide7, 8.8, 1.7, "Premium", "€800", "€350/mese", [
    "Sito completo + landing page",
    "FB + IG + WhatsApp",
    "12-16 post/mese + storie + reels",
    "Sponsorizzate + Google Ads",
    "Report bi-settimanale",
    "Google Business + recensioni"
])

# Budget note
add_rounded_rect(slide7, 0.6, 6.5, 11.8, 0.65, BG_TERTIARY, None)
add_textbox(slide7, 0.9, 6.52, 11.2, 0.6,
            "Nota: Il budget per le sponsorizzate Meta (consigliato €150–300/mese) è a parte e viene pagato "
            "direttamente dal cliente. Ogni proposta può essere personalizzata su misura. Contattaci per un preventivo dedicato.",
            font_size=10, color=TEXT_SECONDARY)

# ════════════════════════════════════════════
# SLIDE 8: Perché Noi
# ════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide8, BG_PRIMARY)

add_textbox(slide8, 0.8, 0.5, 6, 0.6, "Perché Scegliere Noi",
            font_name=FONT_HEADING, font_size=32, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide8, 0.8, 1.1, 1.0)

add_textbox(slide8, 0.8, 1.4, 10, 0.5,
            "Non siamo una grande agenzia impersonale. Siamo un partner locale che conosce il territorio e lavora al tuo fianco.",
            font_size=15, color=TEXT_SECONDARY)

why_cards = [
    ("01", "Conoscenza del Territorio",
     "Viviamo e lavoriamo in Sicilia. Conosciamo le dinamiche, le persone e il modo di comunicare della nostra terra."),
    ("02", "Esperienza Concreta",
     "Un portfolio di progetti reali per attività locali — siti web, gestione social e campagne che hanno portato risultati tangibili."),
    ("03", "Approccio Personale",
     "Un referente unico, sempre disponibile. Nessun call center, nessun ticket. Comunicazione diretta e rapida."),
    ("04", "Risultati Misurabili",
     "Report periodici chiari con numeri reali: persone raggiunte, interazioni, visite al sito. Sai sempre cosa ottieni.")
]

positions = [(0.8, 2.2), (6.3, 2.2), (0.8, 4.5), (6.3, 4.5)]
for (num, title, desc), (x, y) in zip(why_cards, positions):
    card_w = 5.2
    add_rounded_rect(slide8, x, y, card_w, 2.0, BG_CARD, None)
    add_textbox(slide8, x + 0.25, y + 0.15, 1, 0.5, num,
                font_name=FONT_HEADING, font_size=24, color=ACCENT, bold=True)
    add_textbox(slide8, x + 0.25, y + 0.55, card_w - 0.5, 0.4, title,
                font_name=FONT_HEADING, font_size=15, color=ACCENT_LIGHT, bold=True)
    add_textbox(slide8, x + 0.25, y + 1.0, card_w - 0.5, 0.9, desc,
                font_size=12, color=TEXT_SECONDARY)

# ════════════════════════════════════════════
# SLIDE 9: Contatto
# ════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide9, BG_SECONDARY)

add_textbox(slide9, 2, 0.8, 9.3, 0.7, "Prossimi Passi",
            font_name=FONT_HEADING, font_size=36, color=TEXT_PRIMARY,
            bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide9, 6.0, 1.5, 1.3)

add_textbox(slide9, 2, 1.8, 9.3, 0.7,
            "Iniziamo insieme a costruire la tua visibilità.",
            font_name=FONT_HEADING, font_size=22, color=ACCENT_LIGHT,
            italic=True, alignment=PP_ALIGN.CENTER)

# Contact card
card_x, card_y, card_w, card_h = 4.0, 2.8, 5.3, 3.2
add_rounded_rect(slide9, card_x, card_y, card_w, card_h, BG_CARD, ACCENT)

add_textbox(slide9, card_x + 0.4, card_y + 0.25, card_w - 0.8, 0.4,
            "EMC Digital Solutions",
            font_name=FONT_HEADING, font_size=20, color=TEXT_PRIMARY, bold=True)
add_textbox(slide9, card_x + 0.4, card_y + 0.65, card_w - 0.8, 0.3,
            "Full Stack Developers & Graphics & Digital Consultant",
            font_size=12, color=ACCENT)

contact_lines = [
    "+39 333 XXX XXXX",
    "info@emcdigitalsolutions.it",
    "www.emcdigitalsolutions.it"
]
add_multiline_textbox(slide9, card_x + 0.4, card_y + 1.2, card_w - 0.8, 1.8,
                      [(line, {"font_size": 14, "space_after": 12}) for line in contact_lines],
                      color=TEXT_SECONDARY, line_spacing=1.8)

# Footer
add_textbox(slide9, 3.5, 6.5, 6.3, 0.4,
            "Progettato e Sviluppato da EMC Digital Solutions",
            font_size=10, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER)

# ─── Salva ───
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "proposta-commerciale.pptx")
prs.save(output_path)
print(f"PPTX salvato: {output_path}")
